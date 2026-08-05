"""Tests for make_pptx (markdown -> slide deck)."""

from typing import Any

from halia.skills.export import MakePptx, render_markdown_pptx

_MD = """# Term Review

Welcome to the term review.

- Overall average up 6 points
- Attendance steady

---

# Top Performers

| Student | Average |
| --- | --- |
| Chong | 90.0 |
| Aisha | 78.3 |
"""


def _slide_text(slide: Any) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                parts.extend(cell.text for cell in row.cells)
    return " ".join(parts)


def test_splits_into_slides_with_titles() -> None:
    prs = render_markdown_pptx(_MD)
    assert len(prs.slides) == 2
    assert prs.slides[0].shapes.title.text == "Term Review"
    assert prs.slides[1].shapes.title.text == "Top Performers"


def test_bullets_and_table_render() -> None:
    prs = render_markdown_pptx(_MD)
    assert "Overall average up 6 points" in _slide_text(prs.slides[0])
    table_text = _slide_text(prs.slides[1])
    assert "Chong" in table_text and "90.0" in table_text


def test_no_separator_splits_on_headings() -> None:
    md = "# One\n\ntext a\n\n# Two\n\ntext b"
    prs = render_markdown_pptx(md)
    assert len(prs.slides) == 2


def test_make_pptx_writes_only_deck_by_default(tmp_path: Any) -> None:
    out = tmp_path / "deck.pptx"
    result = MakePptx().run({"path": str(out), "content": _MD})
    assert "wrote 2-slide deck" in result
    assert out.read_bytes().startswith(b"PK")  # pptx is a zip
    assert not (tmp_path / "deck.md").exists()  # no sidecar by default


def test_make_pptx_keeps_source_on_request(tmp_path: Any) -> None:
    out = tmp_path / "deck.pptx"
    result = MakePptx().run({"path": str(out), "content": _MD, "keep_source": True})
    assert "markdown source kept" in result
    assert (tmp_path / "deck.md").read_text() == _MD


def test_make_pptx_supports_unicode(tmp_path: Any) -> None:
    # python-pptx is XML/UTF-8 — full Unicode (unlike the latin-1 PDF path)
    out = tmp_path / "u.pptx"
    r = MakePptx().run({"path": str(out), "content": "# Ringkasan\n\n- Kehadiran ↑ 好"})
    assert "wrote" in r
    from pptx import Presentation

    assert Presentation(str(out)).slides[0].shapes.title.text == "Ringkasan"


def test_pptx_embeds_a_native_chart() -> None:
    md = "# Metrics\n\n```chart\ntitle: Scores\nA: 10\nB: 20\nC: 15\n```"
    prs = render_markdown_pptx(md)
    charts = [sh for s in prs.slides for sh in s.shapes if sh.has_chart]
    assert len(charts) == 1  # a real, editable PowerPoint chart (not an image)


def test_make_pptx_validates(tmp_path: Any) -> None:
    assert "content" in MakePptx().run({"path": str(tmp_path / "x.pptx"), "content": " "})


def test_make_pptx_honors_permission_floor(tmp_path: Any) -> None:
    blocked = tmp_path / "credentials.pptx"
    r = MakePptx().run({"path": str(blocked), "content": "# Secret\n\nhi"})
    assert r.startswith("blocked:")
    assert not blocked.exists()


def test_make_pptx_is_dangerous_and_wired() -> None:
    from halia.presets import get_preset
    from halia.skills import available_skills

    assert MakePptx().dangerous is True
    assert "make_pptx" in available_skills()
    assert "make_pptx" in get_preset("education").skills
