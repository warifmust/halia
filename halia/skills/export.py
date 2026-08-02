"""Export skills — render halia's markdown content to printable artifacts.

Per the decided model: **markdown is the master**, and each output format is a direct
*render* from it (never a format→format conversion — that would need LibreOffice).
`make_pdf` is the first render. It keeps the markdown source alongside the PDF, so the
PDF stays a disposable print product and the markdown remains the editable working copy.

Renders a pragmatic markdown subset (headings, paragraphs, bold, bullet/numbered lists,
simple tables, rules) with fpdf2 — lean, no browser, no LibreOffice.
"""

from __future__ import annotations

import math
import re
from pathlib import Path
from typing import Any

from fpdf import FPDF
from fpdf.enums import XPos, YPos

from halia.permissions.guard import PermissionDenied, check_writable
from halia.skills.chart import parse_chart_block

_HEADING_SIZES = {1: 20, 2: 16, 3: 14, 4: 12, 5: 11, 6: 11}
# Core PDF fonts are latin-1 only; map common typographic characters so content
# doesn't crash the renderer. (Full-Unicode output needs a bundled TTF — deferred.)
_REPLACEMENTS = {
    "—": "-", "–": "-", "‘": "'", "’": "'",
    "“": '"', "”": '"', "…": "...", "•": "-", " ": " ",
}


# Bundled Unicode font so PDFs render non-Latin text (Malay, accents, Cyrillic, …).
# If the asset is somehow missing, fall back to a core font + latin-1 sanitisation.
_FONT_DIR = Path(__file__).resolve().parent.parent / "assets" / "fonts"
_FONT_REGULAR = _FONT_DIR / "DejaVuSans.ttf"
_FONT_BOLD = _FONT_DIR / "DejaVuSans-Bold.ttf"
_UNICODE = _FONT_REGULAR.exists() and _FONT_BOLD.exists()
_FONT = "DejaVu" if _UNICODE else "Helvetica"


def _register_fonts(pdf: FPDF) -> None:
    if _UNICODE:
        pdf.add_font("DejaVu", "", str(_FONT_REGULAR))
        pdf.add_font("DejaVu", "B", str(_FONT_BOLD))


def _s(text: str) -> str:
    if _UNICODE:
        return text  # DejaVu handles Unicode natively — no sanitisation needed
    for uni, ascii_ in _REPLACEMENTS.items():
        text = text.replace(uni, ascii_)
    return text.encode("latin-1", "replace").decode("latin-1")


def _truncate(pdf: FPDF, text: str, max_w: float) -> str:
    if pdf.get_string_width(text) <= max_w:
        return text
    while text and pdf.get_string_width(text + "...") > max_w:
        text = text[:-1]
    return text + "..."


def _render_table(pdf: FPDF, block: list[str]) -> None:
    rows: list[list[str]] = []
    for raw in block:
        if re.match(r"^\|[\s:|-]+\|$", raw):  # the |---|---| separator row
            continue
        rows.append([c.strip() for c in raw.strip().strip("|").split("|")])
    if not rows:
        return
    ncols = max(len(r) for r in rows)
    col_w = (pdf.w - pdf.l_margin - pdf.r_margin) / ncols
    for i, row in enumerate(rows):
        cells = row + [""] * (ncols - len(row))
        pdf.set_font(_FONT, "B" if i == 0 else "", 10)
        for cell in cells:
            pdf.cell(col_w, 8, _truncate(pdf, _s(cell), col_w - 2), border=1)
        pdf.ln()
    pdf.ln(2)


def _num(value: float) -> str:
    return str(int(value)) if float(value).is_integer() else f"{value:.1f}"


_CHART_COLORS = [
    (79, 124, 255), (255, 107, 107), (46, 204, 113),
    (241, 196, 15), (155, 89, 182), (230, 126, 34),
]


def _chart_frame(pdf: FPDF, title: str, series: list[tuple[str, list[float]]]) -> float:
    """Common chart setup: page-break, title, legend. Returns the plot-top y."""
    chart_h = 55.0
    if pdf.get_y() + chart_h + 30 > pdf.h - pdf.b_margin:
        pdf.add_page()
    pdf.ln(2)
    if title:
        pdf.set_font(_FONT, "B", 11)
        pdf.multi_cell(pdf.epw, 6, _s(title), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    if len(series) > 1:  # legend row
        pdf.set_font(_FONT, "", 8)
        lx = pdf.l_margin
        y0 = pdf.get_y()
        for si, (name, _) in enumerate(series):
            r, g, b = _CHART_COLORS[si % len(_CHART_COLORS)]
            pdf.set_fill_color(r, g, b)
            pdf.rect(lx, y0 + 0.5, 3, 3, style="F")
            label = name or f"Series {si + 1}"
            pdf.set_xy(lx + 4, y0)
            pdf.cell(pdf.get_string_width(_s(label)) + 2, 4, _s(label))
            lx += 6 + pdf.get_string_width(_s(label)) + 6
        pdf.ln(6)
    return pdf.get_y() + 4


def _draw_chart_pdf(pdf: FPDF, spec: Any) -> None:
    """Draw a chart natively with fpdf2 primitives, dispatched by kind."""
    if spec.kind == "pie":
        _draw_pie_pdf(pdf, spec)
    elif spec.kind == "scatter":
        _draw_scatter_pdf(pdf, spec)
    else:
        _draw_bar_line_pdf(pdf, spec)


def _draw_pie_pdf(pdf: FPDF, spec: Any) -> None:
    chart_h = 62.0
    if pdf.get_y() + chart_h + 30 > pdf.h - pdf.b_margin:
        pdf.add_page()
    pdf.ln(2)
    if spec.title:
        pdf.set_font(_FONT, "B", 11)
        pdf.multi_cell(pdf.epw, 6, _s(spec.title), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    labels = spec.categories
    values = spec.series[0][1] if spec.series else []
    total = sum(values) or 1
    top = pdf.get_y() + 4
    cx, cy, r = pdf.l_margin + 30.0, top + 30.0, 28.0
    angle = -90.0
    for i, value in enumerate(values):
        sweep = value / total * 360
        rr, gg, bb = _CHART_COLORS[i % len(_CHART_COLORS)]
        pdf.set_fill_color(rr, gg, bb)
        if sweep >= 359.999:
            pdf.ellipse(cx - r, cy - r, 2 * r, 2 * r, style="F")
        else:
            steps = max(2, int(sweep / 6))
            poly = [(cx, cy)] + [
                (cx + r * math.cos(math.radians(angle + sweep * s / steps)),
                 cy + r * math.sin(math.radians(angle + sweep * s / steps)))
                for s in range(steps + 1)
            ]
            pdf.polygon(poly, style="F")
        angle += sweep
    pdf.set_font(_FONT, "", 9)
    lx, ly = cx + r + 14, top
    for i, (label, value) in enumerate(zip(labels, values, strict=False)):
        rr, gg, bb = _CHART_COLORS[i % len(_CHART_COLORS)]
        pdf.set_fill_color(rr, gg, bb)
        pdf.rect(lx, ly, 3, 3, style="F")
        pdf.set_xy(lx + 5, ly - 1)
        pdf.cell(70, 4, _s(f"{label} ({value / total * 100:.0f}%)"))
        ly += 6
    pdf.set_y(max(cy + r, ly) + 6)
    pdf.ln(2)


def _draw_scatter_pdf(pdf: FPDF, spec: Any) -> None:
    chart_h = 60.0
    if pdf.get_y() + chart_h + 30 > pdf.h - pdf.b_margin:
        pdf.add_page()
    pdf.ln(2)
    if spec.title:
        pdf.set_font(_FONT, "B", 11)
        pdf.multi_cell(pdf.epw, 6, _s(spec.title), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pts = spec.points
    xs = [p[0] for p in pts] or [0.0]
    ys = [p[1] for p in pts] or [0.0]
    xmin, xmax, ymin, ymax = min(xs), max(xs), min(ys), max(ys)
    xspan, yspan = (xmax - xmin) or 1, (ymax - ymin) or 1
    left, usable = pdf.l_margin, pdf.epw
    top = pdf.get_y() + 4
    baseline = top + chart_h
    pdf.set_draw_color(136, 136, 136)
    pdf.line(left, baseline, left + usable, baseline)
    pdf.line(left, top, left, baseline)
    pdf.set_draw_color(0, 0, 0)
    rr, gg, bb = _CHART_COLORS[0]
    pdf.set_fill_color(rr, gg, bb)
    for x, y in pts:
        px = left + (x - xmin) / xspan * usable
        py = baseline - (y - ymin) / yspan * chart_h
        pdf.ellipse(px - 1, py - 1, 2, 2, style="F")
    if spec.x_label:
        pdf.set_font(_FONT, "", 8)
        pdf.set_xy(left, baseline + 2)
        pdf.cell(usable, 4, _s(spec.x_label), align="C")
    pdf.set_y(baseline + 8)
    pdf.ln(2)


def _draw_bar_line_pdf(pdf: FPDF, spec: Any) -> None:
    """Draw a bar (grouped) or line (multi) chart natively with fpdf2 primitives."""
    kind: str = spec.kind
    title: str = spec.title
    categories: list[str] = spec.categories
    series: list[tuple[str, list[float]]] = spec.series
    chart_h = 55.0
    top = _chart_frame(pdf, title, series)
    left = pdf.l_margin
    usable = pdf.epw
    baseline = top + chart_h
    all_v = [v for _, vals in series for v in vals] or [0.0]
    max_v, min_v = max(all_v), min(0.0, min(all_v))
    span = (max_v - min_v) or 1
    n_cat = len(categories)
    n_ser = len(series)

    def yv(v: float) -> float:
        return baseline - (v - min_v) / span * chart_h

    def val(vals: list[float], ci: int) -> float:
        return vals[ci] if ci < len(vals) else 0.0

    pdf.set_font(_FONT, "", 8)
    if kind == "line":
        step = usable / (n_cat - 1) if n_cat > 1 else 0.0
        for si, (_, vals) in enumerate(series):
            r, g, b = _CHART_COLORS[si % len(_CHART_COLORS)]
            pdf.set_draw_color(r, g, b)
            pdf.set_line_width(0.5)
            pts = [(left + ci * step, yv(val(vals, ci))) for ci in range(n_cat)]
            for (x1, y1), (x2, y2) in zip(pts, pts[1:], strict=False):
                pdf.line(x1, y1, x2, y2)
            if n_ser == 1:
                for (x, y), ci in zip(pts, range(n_cat), strict=True):
                    pdf.set_xy(x - 10, y - 5)
                    pdf.cell(20, 4, _s(_num(val(vals, ci))), align="C")
        pdf.set_draw_color(0, 0, 0)
        pdf.set_line_width(0.2)
        label_x = [left + ci * step for ci in range(n_cat)]
    else:
        slot = usable / max(n_cat, 1)
        group_w = slot * 0.8
        bar_w = group_w / max(n_ser, 1)
        for ci in range(n_cat):
            for si, (_, vals) in enumerate(series):
                r, g, b = _CHART_COLORS[si % len(_CHART_COLORS)]
                pdf.set_fill_color(r, g, b)
                v = val(vals, ci)
                x = left + ci * slot + (slot - group_w) / 2 + si * bar_w
                y = yv(v)
                pdf.rect(x, y, bar_w, baseline - y, style="F")
                if n_ser == 1:
                    pdf.set_xy(x - 2, y - 4)
                    pdf.cell(bar_w + 4, 4, _s(_num(v)), align="C")
        label_x = [left + ci * slot + slot / 2 for ci in range(n_cat)]

    for ci, cx in enumerate(label_x):
        pdf.set_xy(cx - 15, baseline + 1)
        pdf.cell(30, 4, _s(categories[ci][:14]), align="C")
    pdf.line(left, baseline, left + usable, baseline)
    pdf.set_y(baseline + 8)
    pdf.ln(2)


def render_markdown_pdf(content: str) -> FPDF:
    """Render a markdown subset to a clean FPDF document (pure — caller does the I/O)."""
    pdf = FPDF()
    _register_fonts(pdf)
    pdf.set_margins(20, 20, 20)
    pdf.set_auto_page_break(True, margin=18)
    pdf.add_page()

    lines = content.replace("\r\n", "\n").split("\n")
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()

        if stripped.startswith("```chart"):
            i += 1
            chart_lines: list[str] = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                chart_lines.append(lines[i])
                i += 1
            i += 1  # consume the closing fence
            spec = parse_chart_block("\n".join(chart_lines))
            if spec:
                _draw_chart_pdf(pdf, spec)
            continue

        if stripped.startswith("|") and stripped.endswith("|"):
            table_block: list[str] = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_block.append(lines[i].strip())
                i += 1
            _render_table(pdf, table_block)
            continue
        i += 1

        if not stripped:
            pdf.ln(3)
            continue

        heading = re.match(r"^(#{1,6})\s+(.*)", stripped)
        if heading:
            level = len(heading.group(1))
            pdf.ln(2)
            pdf.set_font(_FONT, "B", _HEADING_SIZES.get(level, 12))
            pdf.multi_cell(pdf.epw, _HEADING_SIZES.get(level, 12) * 0.5 + 2,
                           _s(heading.group(2)), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.ln(1)
            continue

        if re.match(r"^(-{3,}|\*{3,}|_{3,})$", stripped):
            pdf.ln(1)
            y = pdf.get_y()
            pdf.line(pdf.l_margin, y, pdf.w - pdf.r_margin, y)
            pdf.ln(3)
            continue

        bullet = re.match(r"^[-*]\s+(.*)", stripped)
        numbered = re.match(r"^(\d+)\.\s+(.*)", stripped)
        pdf.set_font(_FONT, "", 11)
        if bullet:
            pdf.set_x(pdf.l_margin + 6)
            pdf.multi_cell(pdf.epw - 6, 6, _s("- " + bullet.group(1)),
                           markdown=True, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        elif numbered:
            pdf.set_x(pdf.l_margin + 6)
            pdf.multi_cell(pdf.epw - 6, 6, _s(f"{numbered.group(1)}. {numbered.group(2)}"),
                           markdown=True, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        else:
            pdf.multi_cell(pdf.epw, 6, _s(stripped),
                           markdown=True, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.ln(1)
    return pdf


class MakePdf:
    name = "make_pdf"
    description = (
        "Render markdown/text content to a clean, printable PDF (headings, bold, bullet "
        "and numbered lists, simple tables). To embed a bar chart, include a fenced block:\n"
        "```chart\ntype: line\ntitle: My Chart\nLabel A: 12\nLabel B: 8\n```\n"
        "(type is bar or line — line for trends). The editable markdown source is saved "
        "alongside the PDF — edit that and re-render."
    )
    dangerous = True  # writes files
    parameters: dict[str, Any] = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "path": {"type": "string", "description": "Output .pdf file path."},
            "content": {"type": "string", "description": "The markdown/text content."},
            "title": {"type": "string", "description": "Optional title (rendered as a heading)."},
        },
        "required": ["path", "content"],
    }

    def run(self, args: dict[str, Any]) -> str:
        path = args.get("path")
        content = args.get("content")
        title = args.get("title")
        if not isinstance(path, str) or not path.strip():
            return "error: 'path' is required"
        if not isinstance(content, str) or not content.strip():
            return "error: 'content' is required and must be non-empty"
        if isinstance(title, str) and title.strip():
            content = f"# {title.strip()}\n\n{content}"

        pdf_path = Path(path).expanduser()
        md_path = pdf_path.with_suffix(".md")  # the editable master, kept alongside
        try:
            check_writable(pdf_path)
            check_writable(md_path)
        except PermissionDenied as exc:
            return f"blocked: {exc}"

        try:
            render_markdown_pdf(content).output(str(pdf_path))
            md_path.write_text(content, encoding="utf-8")
        except OSError as exc:
            return f"error writing {path}: {exc}"
        return f"wrote PDF to {pdf_path} (editable markdown source kept at {md_path})"


# --- PPTX: markdown -> slides -----------------------------------------------------
# Same master model: `---` on its own line separates slides (else split on # / ##
# headings). Each slide's first heading is its title; bullets/paragraphs/tables become
# the body. python-pptx is XML/UTF-8 based, so full Unicode works here (unlike PDF).


def _strip_md(text: str) -> str:
    return text.replace("**", "").replace("__", "")


def _split_slides(content: str) -> list[str]:
    lines = content.replace("\r\n", "\n").split("\n")
    if any(re.match(r"^-{3,}$", ln.strip()) for ln in lines):
        chunks, cur = [], []  # type: list[str], list[str]
        for ln in lines:
            if re.match(r"^-{3,}$", ln.strip()):
                chunks.append("\n".join(cur))
                cur = []
            else:
                cur.append(ln)
        chunks.append("\n".join(cur))
    else:
        chunks, cur = [], []
        for ln in lines:
            if re.match(r"^#{1,2}\s+", ln) and any(x.strip() for x in cur):
                chunks.append("\n".join(cur))
                cur = [ln]
            else:
                cur.append(ln)
        chunks.append("\n".join(cur))
    return [c for c in chunks if c.strip()] or [content]


def _parse_chunk(chunk: str) -> tuple[str, list[tuple[str, Any]]]:
    lines = chunk.split("\n")
    title = ""
    body: list[str] = []
    for ln in lines:
        head = re.match(r"^#{1,6}\s+(.*)", ln.strip())
        if head and not title:
            title = _strip_md(head.group(1))
        else:
            body.append(ln)

    blocks: list[tuple[str, Any]] = []
    i = 0
    while i < len(body):
        s = body[i].strip()
        if s.startswith("```chart"):
            i += 1
            chart_lines: list[str] = []
            while i < len(body) and not body[i].strip().startswith("```"):
                chart_lines.append(body[i])
                i += 1
            i += 1  # consume closing fence
            spec = parse_chart_block("\n".join(chart_lines))
            if spec:
                blocks.append(("chart", spec))
            continue
        if s.startswith("|") and s.endswith("|"):
            table: list[list[str]] = []
            while i < len(body) and body[i].strip().startswith("|"):
                row = body[i].strip()
                if not re.match(r"^\|[\s:|-]+\|$", row):
                    table.append([_strip_md(c.strip()) for c in row.strip("|").split("|")])
                i += 1
            if table:
                blocks.append(("table", table))
            continue
        i += 1
        if not s:
            continue
        bullet = re.match(r"^[-*]\s+(.*)", s)
        blocks.append(("bullet", _strip_md(bullet.group(1))) if bullet else ("para", _strip_md(s)))
    return title, blocks


def render_markdown_pptx(content: str) -> Any:
    """Render a markdown subset to a python-pptx Presentation (title + bullets + tables)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    for chunk in _split_slides(content):
        title, blocks = _parse_chunk(chunk)
        tables = [b[1] for b in blocks if b[0] == "table"]
        charts = [b[1] for b in blocks if b[0] == "chart"]
        text_blocks = [b for b in blocks if b[0] in ("bullet", "para")]
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        slide.shapes.title.text = title or "Slide"

        top = 1.7
        if text_blocks:
            box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.8), Inches(4))
            frame = box.text_frame
            frame.word_wrap = True
            for idx, (kind, text) in enumerate(text_blocks):
                para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
                para.text = f"• {text}" if kind == "bullet" else text
                para.font.size = Pt(18)
            top += 0.45 * len(text_blocks) + 0.3
        for table in tables:
            top = _add_pptx_table(slide, table, top, Inches, Pt)
        for chart in charts:
            top = _add_pptx_chart(slide, chart, top, Inches)
    return prs


_PPTX_TYPES = {"bar": "COLUMN_CLUSTERED", "line": "LINE_MARKERS", "pie": "PIE"}


def _add_pptx_chart(slide: Any, spec: Any, top: float, Inches: Any) -> float:
    """Add a native, editable PowerPoint chart from a ChartSpec."""
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE

    height = 4.2
    data: Any
    if spec.kind == "scatter":
        data = XyChartData()  # type: ignore[no-untyped-call]
        s = data.add_series(spec.title or "Points")
        for x, y in spec.points:
            s.add_data_point(x, y)
        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER, Inches(0.6), Inches(top), Inches(8.8), Inches(height), data,
        )
    else:
        data = CategoryChartData()  # type: ignore[no-untyped-call]
        data.categories = spec.categories
        for name, values in spec.series:
            data.add_series(name or "Values", values)
        chart_type = getattr(XL_CHART_TYPE, _PPTX_TYPES.get(spec.kind, "COLUMN_CLUSTERED"))
        frame = slide.shapes.add_chart(
            chart_type, Inches(0.6), Inches(top), Inches(8.8), Inches(height), data,
        )
        if len(spec.series) > 1 or spec.kind == "pie":
            frame.chart.has_legend = True
    if spec.title:
        frame.chart.has_title = True
        frame.chart.chart_title.text_frame.text = spec.title
    return top + height + 0.3


def _add_pptx_table(slide: Any, rows: list[list[str]], top: float, Inches: Any, Pt: Any) -> float:
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    height = 0.4 * n_rows
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(top), Inches(8.8),
                                   Inches(height))
    table = shape.table
    for r, row in enumerate(rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.text = row[c] if c < len(row) else ""
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12)
                para.font.bold = r == 0
    return top + height + 0.3


# --- DOCX: markdown -> Word -------------------------------------------------------
# Same master model. Word is XML/UTF-8 so full Unicode works. python-docx has no native
# chart, so a ```chart block renders as a small Label/Value table.


def _add_docx_rich(doc: Any, text: str, style: str | None = None) -> None:
    para = doc.add_paragraph(style=style)
    for part in re.split(r"(\*\*[^*]+\*\*)", text):
        if part.startswith("**") and part.endswith("**"):
            para.add_run(part[2:-2]).bold = True
        elif part:
            para.add_run(part)


def _add_docx_table(doc: Any, raw_lines: list[str]) -> None:
    data: list[list[str]] = []
    for raw in raw_lines:
        if re.match(r"^\|[\s:|-]+\|$", raw):
            continue
        data.append([_strip_md(c.strip()) for c in raw.strip("|").split("|")])
    if not data:
        return
    ncols = max(len(r) for r in data)
    table = doc.add_table(rows=len(data), cols=ncols)
    table.style = "Table Grid"
    for r, row in enumerate(data):
        for c in range(ncols):
            cell = table.cell(r, c)
            cell.text = row[c] if c < len(row) else ""
            if r == 0 and cell.paragraphs[0].runs:
                cell.paragraphs[0].runs[0].bold = True


def render_markdown_docx(content: str) -> Any:
    """Render a markdown subset to a python-docx Document (headings, lists, tables, bold)."""
    from docx import Document

    doc = Document()
    lines = content.replace("\r\n", "\n").split("\n")
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()

        if stripped.startswith("```chart"):
            i += 1
            chart_lines: list[str] = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                chart_lines.append(lines[i])
                i += 1
            i += 1
            spec = parse_chart_block("\n".join(chart_lines))
            if spec:
                if spec.title:
                    doc.add_heading(spec.title, level=3)
                if spec.kind == "scatter":  # python-docx has no chart → a table of (x, y)
                    rows = ["| x | y |"] + [f"| {x} | {y} |" for x, y in spec.points]
                else:
                    rows = ["| Item | " + " | ".join(n or f"Series {i + 1}"
                            for i, (n, _) in enumerate(spec.series)) + " |"] + [
                        "| " + spec.categories[ci] + " | "
                        + " | ".join(str(v[ci]) if ci < len(v) else "" for _, v in spec.series)
                        + " |"
                        for ci in range(len(spec.categories))
                    ]
                _add_docx_table(doc, rows)
            continue

        if stripped.startswith("|") and stripped.endswith("|"):
            table_block: list[str] = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_block.append(lines[i].strip())
                i += 1
            _add_docx_table(doc, table_block)
            continue
        i += 1

        if not stripped:
            continue
        heading = re.match(r"^(#{1,6})\s+(.*)", stripped)
        if heading:
            doc.add_heading(_strip_md(heading.group(2)), level=min(len(heading.group(1)), 4))
            continue
        if re.match(r"^(-{3,}|\*{3,}|_{3,})$", stripped):
            doc.add_paragraph()
            continue
        bullet = re.match(r"^[-*]\s+(.*)", stripped)
        numbered = re.match(r"^(\d+)\.\s+(.*)", stripped)
        if bullet:
            _add_docx_rich(doc, bullet.group(1), "List Bullet")
        elif numbered:
            _add_docx_rich(doc, numbered.group(2), "List Number")
        else:
            _add_docx_rich(doc, stripped)
    return doc


class MakeDocx:
    name = "make_docx"
    description = (
        "Render markdown/text content to an editable Word (.docx) document (headings, bold, "
        "bullet and numbered lists, simple tables). Full Unicode. The editable markdown "
        "source is saved alongside — edit that and re-render."
    )
    dangerous = True  # writes files
    parameters: dict[str, Any] = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "path": {"type": "string", "description": "Output .docx file path."},
            "content": {"type": "string", "description": "The markdown/text content."},
            "title": {"type": "string", "description": "Optional title (rendered as a heading)."},
        },
        "required": ["path", "content"],
    }

    def run(self, args: dict[str, Any]) -> str:
        path = args.get("path")
        content = args.get("content")
        title = args.get("title")
        if not isinstance(path, str) or not path.strip():
            return "error: 'path' is required"
        if not isinstance(content, str) or not content.strip():
            return "error: 'content' is required and must be non-empty"
        if isinstance(title, str) and title.strip():
            content = f"# {title.strip()}\n\n{content}"

        docx_path = Path(path).expanduser()
        md_path = docx_path.with_suffix(".md")
        try:
            check_writable(docx_path)
            check_writable(md_path)
        except PermissionDenied as exc:
            return f"blocked: {exc}"

        try:
            render_markdown_docx(content).save(str(docx_path))
            md_path.write_text(content, encoding="utf-8")
        except OSError as exc:
            return f"error writing {path}: {exc}"
        return f"wrote Word document to {docx_path} (editable markdown source kept at {md_path})"


class MakePptx:
    name = "make_pptx"
    description = (
        "Render markdown content to a PowerPoint (.pptx) deck of structured content slides "
        "(title + bullets + simple tables + native editable charts). Use '---' on its own "
        "line to separate slides. Embed a chart with a fenced block:\n"
        "```chart\ntype: line\ntitle: My Chart\nLabel A: 12\nLabel B: 8\n```\n"
        "(type is bar or line). Produces clean CONTENT and arrangement; the user styles the "
        "design in PowerPoint. "
        "The editable markdown source is saved alongside."
    )
    dangerous = True  # writes files
    parameters: dict[str, Any] = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "path": {"type": "string", "description": "Output .pptx file path."},
            "content": {
                "type": "string",
                "description": "Markdown content; '---' on its own line starts a new slide.",
            },
        },
        "required": ["path", "content"],
    }

    def run(self, args: dict[str, Any]) -> str:
        path = args.get("path")
        content = args.get("content")
        if not isinstance(path, str) or not path.strip():
            return "error: 'path' is required"
        if not isinstance(content, str) or not content.strip():
            return "error: 'content' is required and must be non-empty"

        pptx_path = Path(path).expanduser()
        md_path = pptx_path.with_suffix(".md")
        try:
            check_writable(pptx_path)
            check_writable(md_path)
        except PermissionDenied as exc:
            return f"blocked: {exc}"

        deck = render_markdown_pptx(content)
        n = len(deck.slides)
        try:
            deck.save(str(pptx_path))
            md_path.write_text(content, encoding="utf-8")
        except OSError as exc:
            return f"error writing {path}: {exc}"
        return f"wrote {n}-slide deck to {pptx_path} (editable markdown source kept at {md_path})"
